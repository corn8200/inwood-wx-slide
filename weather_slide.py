#!/usr/bin/env python3
"""
Daily weather-heat-stress slide for Inwood, WV

• Pulls 10-day forecast from Open-Meteo (no key required).
• Computes NWS heat-index for the next 24 h.
• Maps HI to your plant’s heat-stress guidance.
• Builds a one-slide PowerPoint.
• E-mails it via SendGrid.

Required libs: python-dotenv, requests, python-pptx, sendgrid
"""

import os, base64, requests
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv
from sendgrid import SendGridAPIClient
from sendgrid.helpers.mail import Mail, Attachment, FileContent, FileName, \
                                   FileType, Disposition

# ── 1. configuration ──────────────────────────────────────────────────────────
load_dotenv()                               # reads .env in project root

LAT  = os.getenv("LATITUDE")                # 39.36  (Inwood)
LON  = os.getenv("LONGITUDE")               # -78.05
TZ   = os.getenv("TIMEZONE")                # America/New_York
SG   = os.getenv("SENDGRID_API_KEY")
FROM = os.getenv("EMAIL_FROM")              # alerts@jcornelius.net
TO   = [e.strip() for e in os.getenv("EMAIL_TO").split(",")]

TODAY = datetime.now().astimezone().strftime("%Y-%m-%d")
SLIDE = f"WX_{TODAY}.pptx"

# ── 2. get forecast ───────────────────────────────────────────────────────────
url = (
    "https://api.open-meteo.com/v1/forecast"
    f"?latitude={LAT}&longitude={LON}"
    "&hourly=temperature_2m,relative_humidity_2m"
    "&daily=temperature_2m_max,relative_humidity_2m_max"
    "&forecast_days=10"
    f"&timezone={TZ}"
)
wx = requests.get(url, timeout=15).json()

temps_c  = wx["hourly"]["temperature_2m"][:24]        # next 24 h
rh_pct   = wx["hourly"]["relative_humidity_2m"][:24]

def c2f(c): return c*9/5 + 32
def heat_idx_f(t_f, rh):                              # Rothfusz regression
    c1,c2,c3,c4,c5,c6,c7,c8,c9 = (
        -42.379, 2.04901523, 10.14333127, -0.22475541,
        -6.83783e-3, -5.481717e-2, 1.22874e-3,
        8.5282e-4, -1.99e-6
    )
    return (c1 + c2*t_f + c3*rh + c4*t_f*rh + c5*t_f**2 +
            c6*rh**2 + c7*t_f**2*rh + c8*t_f*rh**2 + c9*t_f**2*rh**2)

hi_max = round(max(heat_idx_f(c2f(t), rh) for t, rh in zip(temps_c, rh_pct)))

# ── 3. map to plant heat-stress table ─────────────────────────────────────────
table = [
    ( 80,  90, "Caution",          "30 min", "1 cup/20 min", "Normal",    "Periodic"),
    ( 91, 103, "Extreme Caution",  "15 min", "1 cup/15 min", "30-40/10",  "1× hr"),
    (104,124, "Danger",            "10 min", "1 cup/10 min", "20-30/10",  "2× hr"),
    (125,999, "Extreme Danger",     "0",     "1 cup/10 min*", "10-20/10", "4× hr"),
]
for lo, hi, cat, work_max, hydration, work_rest, supv in table:
    if lo <= hi_max <= hi:
        category = cat
        break

# ── 4. build slide ───────────────────────────────────────────────────────────
prs   = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])     # blank

title = slide.shapes.title
title.text = f"Inwood Weather — {TODAY}"

body = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(9), Inches(2.2))
tf   = body.text_frame
p    = tf.paragraphs[0]
p.text = (f"24-h Max Heat Index: {hi_max} °F  ({category})\n"
          "See table below for work/rest & hydration guidance.")
p.font.size = Pt(18)

days = wx["daily"]["time"]
tmax = wx["daily"]["temperature_2m_max"]
rmax = wx["daily"]["relative_humidity_2m_max"]

rows, cols = 11, 3
tbl = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(3.0),
                             Inches(9),   Inches(2.5)).table
tbl.columns[0].width = tbl.columns[1].width = tbl.columns[2].width = Inches(3)

tbl.cell(0,0).text, tbl.cell(0,1).text, tbl.cell(0,2).text = "Date", "Tmax °C", "RH %"

for i in range(10):
    tbl.cell(i+1,0).text = days[i]
    tbl.cell(i+1,1).text = str(tmax[i])
    tbl.cell(i+1,2).text = str(rmax[i])

prs.save(SLIDE)

# ── 5. e-mail slide ──────────────────────────────────────────────────────────
sg  = SendGridAPIClient(SG)
with open(SLIDE, "rb") as f:
    encoded = base64.b64encode(f.read()).decode()

att = Attachment(
    FileContent(encoded),
    FileName(SLIDE),
    FileType("application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    Disposition("attachment")
)

mail = Mail(
    from_email=FROM,
    to_emails=TO,
    subject=f"John’s WX Slide {TODAY}",
    plain_text_content="Daily weather & heat-stress slide attached.",
    attachments=[att],
)
sg.send(mail)
print("Slide emailed successfully.")
